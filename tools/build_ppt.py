"""
Build the PPT for "包裹发货后 常出现的质检不合格的原因".
Style: orange + white, two-column layout matching the example image.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

# ---------- Colors ----------
ORANGE        = RGBColor(0xF4, 0x7B, 0x2A)   # primary
ORANGE_DEEP   = RGBColor(0xE8, 0x65, 0x14)
ORANGE_LIGHT  = RGBColor(0xFF, 0xF1, 0xE6)
MAGENTA       = RGBColor(0xE8, 0x35, 0x8C)   # accent
GRAY_LIGHT    = RGBColor(0xF5, 0xF5, 0xF5)
GRAY_BORDER   = RGBColor(0xE5, 0xE5, 0xE5)
GRAY_TEXT     = RGBColor(0x59, 0x59, 0x59)
DARK_TEXT     = RGBColor(0x2B, 0x2B, 0x2B)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)

# ---------- Geometry (16:9 widescreen 13.333" x 7.5") ----------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MEDIA_DIR = "tools/extracted/media"
OUT = "包裹发货后常见质检不合格原因.pptx"


def add_rect(slide, x, y, w, h, fill_color, line_color=None, shape=MSO_SHAPE.RECTANGLE):
    shp = slide.shapes.add_shape(shape, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(0.5)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=DARK_TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="微软雅黑"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_circle_number(slide, x, y, diameter, number, fill=ORANGE, text_color=WHITE):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, diameter, diameter)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    tf = shp.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(number)
    run.font.name = "微软雅黑"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = text_color
    return shp


def add_square_marker(slide, x, y, size, color=MAGENTA):
    return add_rect(slide, x, y, size, size, color)


def add_header(slide, title, subtitle=None):
    """Top orange bar + optional subtitle bar."""
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.85), ORANGE)
    add_text(slide, Inches(0.45), Inches(0.12), Inches(12), Inches(0.6),
             title, size=26, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE, font="微软雅黑")
    if subtitle:
        add_rect(slide, Inches(0), Inches(0.85), SLIDE_W, Inches(0.45), GRAY_LIGHT)
        add_text(slide, Inches(0.45), Inches(0.88), Inches(12), Inches(0.4),
                 subtitle, size=12, color=GRAY_TEXT,
                 anchor=MSO_ANCHOR.MIDDLE, font="微软雅黑")


def add_panel_header(slide, x, y, w, label, fill=ORANGE):
    """Orange/magenta header bar for a column panel."""
    h = Inches(0.55)
    add_rect(slide, x, y, w, h, fill)
    add_text(slide, x + Inches(0.25), y, w - Inches(0.5), h,
             label, size=15, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE, font="微软雅黑")


def add_panel_body(slide, x, y, w, h):
    """White panel body with light gray border."""
    add_rect(slide, x, y, w, h, WHITE, line_color=GRAY_BORDER)


def add_image_fit(slide, path, x, y, max_w, max_h):
    """Place image centered in (x,y,max_w,max_h) keeping aspect."""
    img = Image.open(path)
    iw, ih = img.size
    ratio = min(max_w / iw, max_h / ih) * 0.92  # small padding
    w = int(iw * ratio)
    h = int(ih * ratio)
    cx = x + (max_w - w) / 2
    cy = y + (max_h - h) / 2
    slide.shapes.add_picture(path, Emu(int(cx)), Emu(int(cy)), Emu(w), Emu(h))


# =====================================================================
# Build presentation
# =====================================================================
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]


# ---------------------------------------------------------------------
# Slide 1: Cover
# ---------------------------------------------------------------------
def slide_cover():
    s = prs.slides.add_slide(blank_layout)
    # left orange block
    add_rect(s, Inches(0), Inches(0), Inches(5.2), SLIDE_H, ORANGE)
    # decorative deep-orange band
    add_rect(s, Inches(0), Inches(6.2), Inches(5.2), Inches(0.18), ORANGE_DEEP)
    # tag
    add_rect(s, Inches(0.7), Inches(1.4), Inches(1.6), Inches(0.4), WHITE)
    add_text(s, Inches(0.7), Inches(1.4), Inches(1.6), Inches(0.4),
             "质检指南", size=12, bold=True, color=ORANGE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # main title
    add_text(s, Inches(0.7), Inches(2.2), Inches(4.5), Inches(1.2),
             "包裹发货后", size=40, bold=True, color=WHITE)
    add_text(s, Inches(0.7), Inches(3.05), Inches(4.5), Inches(1.5),
             "常出现的质检\n不合格原因", size=40, bold=True, color=WHITE)
    add_text(s, Inches(0.7), Inches(5.4), Inches(4.5), Inches(0.5),
             "Quality Inspection · Failure Analysis", size=12,
             color=ORANGE_LIGHT, font="Arial")

    # right white area – decorative
    add_text(s, Inches(6.0), Inches(2.4), Inches(7), Inches(0.6),
             "7  类  常  见  问  题  全  解  析", size=22, bold=True,
             color=ORANGE)
    add_rect(s, Inches(6.0), Inches(3.15), Inches(0.7), Inches(0.06), ORANGE)
    add_text(s, Inches(6.0), Inches(3.4), Inches(7), Inches(2.5),
             "图物不符 · 翻译错误 · 色差 · 克重\n标题修改 · 洗水唛 · 申诉流程",
             size=16, color=GRAY_TEXT)
    add_text(s, Inches(6.0), Inches(6.4), Inches(7), Inches(0.5),
             "面向 TEMU 跨境电商商家的实操手册",
             size=11, color=GRAY_TEXT)


# ---------------------------------------------------------------------
# Slide 2: Overview / TOC
# ---------------------------------------------------------------------
def slide_overview():
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "目  录  &  问  题  全  景",
               "本手册将围绕以下 7 类质检不合格场景，逐一拆解原因与应对方案")

    items = [
        ("1", "图物不符（下装）", "裤装口袋、纽扣、拉链等细节差异"),
        ("2", "图物不符（上装）", "领口条纹、内衬颜色等差异"),
        ("3", "质检台翻译问题", "标题被错误翻译，引发误判"),
        ("4", "色差问题",       "灯光偏暗导致的颜色判定偏差"),
        ("5", "克重问题",       "实物与样品留底克重不一致"),
        ("6", "标题修改",       "JIT/定制单与普通备货单处理差异"),
        ("7", "洗水唛问题",     "成分、八国语言、二维码等合规要点"),
    ]
    # 2 columns x up to 4 rows
    col_w = Inches(6.0)
    row_h = Inches(1.15)
    start_x = [Inches(0.6), Inches(6.8)]
    start_y = Inches(1.7)
    for i, (num, title, desc) in enumerate(items):
        col = i // 4
        row = i % 4
        x = start_x[col]
        y = start_y + row * (row_h + Inches(0.1))
        # card background
        add_rect(s, x, y, col_w, row_h, GRAY_LIGHT)
        # left orange stripe
        add_rect(s, x, y, Inches(0.12), row_h, ORANGE)
        # number circle
        add_circle_number(s, x + Inches(0.35), y + Inches(0.27), Inches(0.6), num)
        # title
        add_text(s, x + Inches(1.15), y + Inches(0.18), col_w - Inches(1.4),
                 Inches(0.45), title, size=16, bold=True, color=DARK_TEXT)
        # desc
        add_text(s, x + Inches(1.15), y + Inches(0.62), col_w - Inches(1.4),
                 Inches(0.45), desc, size=11, color=GRAY_TEXT)


# ---------------------------------------------------------------------
# Slide helper: standard "issue" page with right-side image(s)
# ---------------------------------------------------------------------
def issue_slide_with_images(num, title, subtitle, points, image_paths,
                            panel_label="案例图示", panel_fill=MAGENTA,
                            note=None):
    s = prs.slides.add_slide(blank_layout)
    add_header(s, f"{num}.  {title}", subtitle)

    # left column: definition + points
    left_x = Inches(0.45)
    left_y = Inches(1.55)
    left_w = Inches(6.2)
    left_h = Inches(5.6)
    # left header
    add_panel_header(s, left_x, left_y, left_w, "原  因  与  关  键  点", fill=ORANGE)
    add_panel_body(s, left_x, left_y + Inches(0.55), left_w, left_h - Inches(0.55))

    inner_x = left_x + Inches(0.4)
    inner_y = left_y + Inches(0.85)
    for i, (head, body) in enumerate(points):
        y = inner_y + Inches(i * 1.05)
        add_circle_number(s, inner_x, y, Inches(0.55), str(i + 1))
        add_text(s, inner_x + Inches(0.85), y - Inches(0.02), left_w - Inches(1.6),
                 Inches(0.4), head, size=14, bold=True, color=DARK_TEXT)
        add_text(s, inner_x + Inches(0.85), y + Inches(0.38), left_w - Inches(1.6),
                 Inches(0.6), body, size=11, color=GRAY_TEXT)

    if note:
        add_rect(s, left_x + Inches(0.3), left_y + left_h - Inches(0.95),
                 left_w - Inches(0.6), Inches(0.7), ORANGE_LIGHT)
        add_text(s, left_x + Inches(0.55), left_y + left_h - Inches(0.85),
                 left_w - Inches(1.0), Inches(0.55),
                 f"💡  {note}", size=11, bold=True, color=ORANGE_DEEP,
                 anchor=MSO_ANCHOR.MIDDLE)

    # right column: images
    right_x = Inches(6.85)
    right_y = Inches(1.55)
    right_w = Inches(6.05)
    right_h = Inches(5.6)
    add_panel_header(s, right_x, right_y, right_w, panel_label, fill=panel_fill)
    add_panel_body(s, right_x, right_y + Inches(0.55), right_w, right_h - Inches(0.55))

    if not image_paths:
        return s

    # Layout images: stack vertically up to 3
    n = len(image_paths)
    slot_top = right_y + Inches(0.75)
    slot_h_total = right_h - Inches(0.85)
    if n == 1:
        add_image_fit(s, image_paths[0], right_x + Inches(0.25), slot_top,
                      right_w - Inches(0.5), slot_h_total)
    else:
        slot_h = slot_h_total / n
        for i, p in enumerate(image_paths):
            y = slot_top + i * slot_h
            add_image_fit(s, p, right_x + Inches(0.25), y,
                          right_w - Inches(0.5), slot_h - Inches(0.15))
    return s


# ---------------------------------------------------------------------
# Slide 3: Issue 1 — Pants mismatch
# ---------------------------------------------------------------------
def slide_issue_pants():
    issue_slide_with_images(
        num="01",
        title="图物不符 · 下装类（裤装）",
        subtitle="实际商品与主图/轮播图存在差异，是质检不合格最常见的原因",
        points=[
            ("裤子口袋背面的设计差异", "实物口袋形状/缝线与主图不符，需核对每一处细节"),
            ("裤子口袋纽扣颗数不一致", "纽扣数量是高频差异点，需逐一比对主图与样衣"),
            ("裤子拉链处颜色/设计差异", "轮播图与抽检实物的拉链颜色明显不同，必须更新主图后才能允收"),
        ],
        image_paths=[
            f"{MEDIA_DIR}/image1.jpeg",
            f"{MEDIA_DIR}/image2.png",
            f"{MEDIA_DIR}/image3.jpeg",
        ],
        panel_label="实  物  对  比  图",
        note="主图更新成功后，质检方可允收"
    )


# ---------------------------------------------------------------------
# Slide 4: Issue 2 — Tops mismatch
# ---------------------------------------------------------------------
def slide_issue_tops():
    issue_slide_with_images(
        num="02",
        title="图物不符 · 上装类（POLO/T恤/外套）",
        subtitle="领口、内衬等细节同样是高频差异点",
        points=[
            ("领口条纹色条存在差异", "POLO/T 恤类商品，领口印花条颜色与轮播图不一致"),
            ("外套内衬颜色存在差异", "实物外套内衬为灰色，但主图显示为黑色，需立即更新"),
        ],
        image_paths=[
            f"{MEDIA_DIR}/image4.png",
            f"{MEDIA_DIR}/image5.png",
        ],
        panel_label="实  物  对  比  图",
        note="积极、细致地调整主图与轮播图，做到「图物一致」"
    )


# ---------------------------------------------------------------------
# Slide 5: Issue 3 — Translation error
# ---------------------------------------------------------------------
def slide_issue_translation():
    issue_slide_with_images(
        num="03",
        title="质检台翻译问题",
        subtitle="标题被翻译错误，导致质检判定与实物不符",
        points=[
            ("典型场景", "质检不合格原因写「标题为两件套，实物是一条短裤」"),
            ("商家自查", "在标题中确认是否真的出现了「两件套」描述"),
            ("处理方式", "若标题中并无「两件套」字样 → 属翻译错误，及时找买手发起申诉"),
        ],
        image_paths=[
            f"{MEDIA_DIR}/image6.png",
            f"{MEDIA_DIR}/image7.png",
        ],
        panel_label="质  检  截  图  &  自  查",
        note="确认非自身错误后立即通知买手发起申诉"
    )


# ---------------------------------------------------------------------
# Slide 6: Issue 4 — Color variance
# ---------------------------------------------------------------------
def slide_issue_color():
    issue_slide_with_images(
        num="04",
        title="色差问题",
        subtitle="质检场地灯光偏暗等原因，可能导致色差判定不予允收",
        points=[
            ("可能成因", "质检场地灯光偏暗 / 角度偏差 → 误判为色差较大"),
            ("商家应对", "提供自然光线下拍摄的图片或视频佐证"),
            ("举证标准", "证明商品实物颜色与轮播图、系统申报颜色一致或相近"),
        ],
        image_paths=[
            f"{MEDIA_DIR}/image8.png",
            f"{MEDIA_DIR}/image9.png",
            f"{MEDIA_DIR}/image10.png",
        ],
        panel_label="质  检  图  vs  佐  证  图",
        note="自然光下的清晰对比图/视频是最有力的申诉证据"
    )


# ---------------------------------------------------------------------
# Slide 7: Issue 5 — Weight
# ---------------------------------------------------------------------
def slide_issue_weight():
    issue_slide_with_images(
        num="05",
        title="克重问题",
        subtitle="实物克重与样品留底不一致，可能未达最低允收标准",
        points=[
            ("问题表现", "抽样质检发现实物重量与留底样品克重不一致"),
            ("风险临界", "实物比样品轻太多 → 可能不到最低允收标准而被驳回"),
            ("商家应对", "提供自己样品的克重图发给买手"),
            ("申诉路径", "买手据此计算是否符合最低允收标准 → 符合即可申诉"),
        ],
        image_paths=[
            f"{MEDIA_DIR}/image11.jpeg",
            f"{MEDIA_DIR}/image12.jpeg",
        ],
        panel_label="样  品  克  重  实  拍",
        note="保留每批次样品克重照片，是规避此风险的关键"
    )


# ---------------------------------------------------------------------
# Slide 8: Issue 6 — Title modification (no images, two-column)
# ---------------------------------------------------------------------
def slide_issue_title():
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "06.  关  于  标  题  修  改",
               "JIT / 定制 / 加急 单不接受抽检不合格后再修改返单信息，需区分情况处理")

    # Top callout
    callout_y = Inches(1.55)
    add_rect(s, Inches(0.45), callout_y, Inches(12.45), Inches(0.7), ORANGE_LIGHT)
    add_rect(s, Inches(0.45), callout_y, Inches(0.12), Inches(0.7), ORANGE)
    add_text(s, Inches(0.75), callout_y, Inches(12), Inches(0.7),
             "⚠  JIT / 定制单 / 加急：抽检不合格后不接受返单信息修改 → "
             "立即下架商品 + 申请缺货 + 信息更正后再上架，是降低货损的最佳方式",
             size=12, bold=True, color=ORANGE_DEEP, anchor=MSO_ANCHOR.MIDDLE)

    # Two cards
    card_y = Inches(2.55)
    card_h = Inches(4.5)

    # Card 1
    x1 = Inches(0.45)
    w1 = Inches(6.2)
    add_panel_header(s, x1, card_y, w1, "情  况  ①  商家自定义标题错误", fill=ORANGE)
    add_panel_body(s, x1, card_y + Inches(0.55), w1, card_h - Inches(0.55))
    add_text(s, x1 + Inches(0.4), card_y + Inches(0.85), w1 - Inches(0.8),
             Inches(0.5), "商家自己定义错了标题", size=13, bold=True, color=DARK_TEXT)
    items1 = [
        ("发起方", "商家自行更改货品标题；英文标题需要找买手处理"),
        ("普通备货单 · 首单", "可以尝试发起申诉"),
        ("普通备货单 · 非首单", "需在所有信息更新后再发起返单"),
        ("JIT / 定制 / 加急", "立即申请缺货 → 减少货损（质检端信息已固化）"),
    ]
    for i, (h, b) in enumerate(items1):
        y = card_y + Inches(1.45) + Inches(i * 0.78)
        add_square_marker(s, x1 + Inches(0.4), y + Inches(0.08), Inches(0.18), color=ORANGE)
        add_text(s, x1 + Inches(0.7), y, w1 - Inches(1.0), Inches(0.3),
                 h, size=12, bold=True, color=DARK_TEXT)
        add_text(s, x1 + Inches(0.7), y + Inches(0.32), w1 - Inches(1.0), Inches(0.4),
                 b, size=10.5, color=GRAY_TEXT)

    # Card 2
    x2 = Inches(6.85)
    w2 = Inches(6.05)
    add_panel_header(s, x2, card_y, w2, "情  况  ②  图灵翻译出错", fill=MAGENTA)
    add_panel_body(s, x2, card_y + Inches(0.55), w2, card_h - Inches(0.55))
    add_text(s, x2 + Inches(0.4), card_y + Inches(0.85), w2 - Inches(0.8), Inches(0.5),
             "商家自定义标题正确，翻译后出现错误", size=13, bold=True, color=DARK_TEXT)
    items2 = [
        ("处理方", "找买手 + 相关负责人员，发起对外语标题的更改申请"),
        ("适用范围", "所有标题翻译错误的情况"),
        ("注意事项", "保留对照截图（中文 vs 翻译后），便于审核"),
        ("时效要求", "尽快发起，避免下一批继续因翻译问题失败"),
    ]
    for i, (h, b) in enumerate(items2):
        y = card_y + Inches(1.45) + Inches(i * 0.78)
        add_square_marker(s, x2 + Inches(0.4), y + Inches(0.08), Inches(0.18), color=MAGENTA)
        add_text(s, x2 + Inches(0.7), y, w2 - Inches(1.0), Inches(0.3),
                 h, size=12, bold=True, color=DARK_TEXT)
        add_text(s, x2 + Inches(0.7), y + Inches(0.32), w2 - Inches(1.0), Inches(0.4),
                 b, size=10.5, color=GRAY_TEXT)


# ---------------------------------------------------------------------
# Slide 9: Issue 7 — Care label
# ---------------------------------------------------------------------
def slide_issue_care_label():
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "07.  洗  水  唛  问  题",
               "成分申报、八国语言、欧盟进口商电子地址、二维码等多维合规检查")

    # Two columns: 5 issues split 3 + 2
    left_items = [
        ("成分不一致", "水洗唛上的成分与后台申报的成分不一致 → 核对一致后截图发买手申诉"),
        ("缺少八国语言", "未车缝八国语言洗水唛 → 核对洗水唛是否包含八国语言，尽早发买手申诉"),
        ("套装未车缝洗水唛", "上装与裤装均需车缝洗水唛 → 若两件都已车缝，找买手申诉"),
    ]
    right_items = [
        ("缺欧盟进口商电子地址", "核对是否有 EU Importer Electronic Address 字样"),
        ("二维码扫码异常", "扫描洗水唛上的二维码，确认是否存在质检提到的问题"),
    ]

    left_x = Inches(0.45); left_w = Inches(6.2)
    right_x = Inches(6.85); right_w = Inches(6.05)
    panel_y = Inches(1.55); panel_h = Inches(5.6)

    add_panel_header(s, left_x, panel_y, left_w, "成  分  &  语  言  &  车  缝", fill=ORANGE)
    add_panel_body(s, left_x, panel_y + Inches(0.55), left_w, panel_h - Inches(0.55))

    add_panel_header(s, right_x, panel_y, right_w, "合  规  &  二  维  码", fill=MAGENTA)
    add_panel_body(s, right_x, panel_y + Inches(0.55), right_w, panel_h - Inches(0.55))

    def render_items(x, w, items, num_color, start_num):
        for i, (h, b) in enumerate(items):
            y = panel_y + Inches(0.95) + Inches(i * 1.55)
            add_circle_number(s, x + Inches(0.4), y, Inches(0.6),
                              str(start_num + i), fill=num_color)
            add_text(s, x + Inches(1.2), y, w - Inches(1.5), Inches(0.5),
                     h, size=15, bold=True, color=DARK_TEXT)
            add_text(s, x + Inches(1.2), y + Inches(0.5), w - Inches(1.5), Inches(0.95),
                     b, size=11, color=GRAY_TEXT)

    render_items(left_x, left_w, left_items, ORANGE, 1)
    render_items(right_x, right_w, right_items, MAGENTA, 4)


# ---------------------------------------------------------------------
# Slide 10: Action plan / appeal flow summary
# ---------------------------------------------------------------------
def slide_action_plan():
    s = prs.slides.add_slide(blank_layout)
    add_header(s, "通  用  应  对  策  略  &  申  诉  流  程",
               "面对任何一类质检不合格，均可遵循以下闭环流程")

    # 5-step horizontal flow
    steps = [
        ("识别原因",   "对照质检台反馈，定位属于 7 类问题中的哪一种"),
        ("自查核对",   "比对主图/轮播图/标题/洗水唛/克重等证据"),
        ("准备证据",   "自然光照片、克重图、对照截图、洗水唛特写"),
        ("发起申诉",   "联系买手提交证据，区分首单/非首单/JIT 处理"),
        ("源头改进",   "更新主图、修正标题翻译、规范洗水唛、补留克重"),
    ]
    n = len(steps)
    margin = Inches(0.5)
    available = SLIDE_W - 2 * margin
    gap = Inches(0.15)
    card_w = (available - gap * (n - 1)) / n
    card_y = Inches(1.7)
    card_h = Inches(2.6)

    for i, (h, b) in enumerate(steps):
        x = margin + i * (card_w + gap)
        # card
        add_rect(s, x, card_y, card_w, card_h, GRAY_LIGHT)
        # top stripe
        add_rect(s, x, card_y, card_w, Inches(0.2), ORANGE)
        # number circle (overlapping top)
        add_circle_number(s, x + (card_w - Inches(0.85)) / 2,
                          card_y + Inches(0.45), Inches(0.85),
                          str(i + 1))
        add_text(s, x + Inches(0.1), card_y + Inches(1.45),
                 card_w - Inches(0.2), Inches(0.45),
                 h, size=15, bold=True, color=DARK_TEXT,
                 align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.15), card_y + Inches(1.92),
                 card_w - Inches(0.3), Inches(0.65),
                 b, size=10.5, color=GRAY_TEXT, align=PP_ALIGN.CENTER)
        # arrow
        if i < n - 1:
            ax = x + card_w + Inches(0.005)
            ay = card_y + card_h / 2 - Inches(0.1)
            arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,
                                       ax, ay, Inches(0.14), Inches(0.2))
            arrow.rotation = 30
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ORANGE
            arrow.line.fill.background()

    # Bottom: golden rules box
    box_y = Inches(4.7)
    add_rect(s, Inches(0.5), box_y, SLIDE_W - Inches(1.0), Inches(2.4), WHITE,
             line_color=GRAY_BORDER)
    add_rect(s, Inches(0.5), box_y, Inches(0.12), Inches(2.4), ORANGE)
    add_text(s, Inches(0.85), box_y + Inches(0.2), Inches(11), Inches(0.5),
             "🔑  核心原则",
             size=16, bold=True, color=ORANGE)
    rules = [
        "图物一致：主图/轮播图必须与实物完全一致，包括口袋、纽扣、拉链、领口、内衬等细节",
        "证据先行：色差需自然光对照图，克重需样品留底实拍图，翻译错误需中外文对照截图",
        "时效优先：JIT / 定制 / 加急 单立即申请缺货，不要尝试返单修改，避免货损",
        "源头治理：每次申诉成功后，同步更新前端信息（主图、标题、洗水唛），杜绝再次发生",
    ]
    for i, r in enumerate(rules):
        y = box_y + Inches(0.7 + i * 0.42)
        add_square_marker(s, Inches(0.95), y + Inches(0.08), Inches(0.16), color=MAGENTA)
        add_text(s, Inches(1.25), y, Inches(11), Inches(0.4),
                 r, size=12, color=DARK_TEXT)


# ---------------------------------------------------------------------
# Slide 11: Closing
# ---------------------------------------------------------------------
def slide_closing():
    s = prs.slides.add_slide(blank_layout)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, ORANGE)
    add_rect(s, Inches(0), Inches(3.55), SLIDE_W, Inches(0.05), WHITE)
    add_text(s, Inches(0), Inches(2.4), SLIDE_W, Inches(1.2),
             "THANKS", size=72, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, font="Arial")
    add_text(s, Inches(0), Inches(3.85), SLIDE_W, Inches(0.6),
             "图  物  一  致  ·  证  据  完  备  ·  时  效  优  先  ·  源  头  治  理",
             size=18, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0), Inches(5.2), SLIDE_W, Inches(0.4),
             "—  让每一个包裹都顺利通过质检  —",
             size=13, color=ORANGE_LIGHT, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------
slide_cover()
slide_overview()
slide_issue_pants()
slide_issue_tops()
slide_issue_translation()
slide_issue_color()
slide_issue_weight()
slide_issue_title()
slide_issue_care_label()
slide_action_plan()
slide_closing()

prs.save(OUT)
print(f"OK -> {OUT}  ({os.path.getsize(OUT) / 1024:.1f} KB)")
